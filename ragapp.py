import streamlit as st
import os
import shutil
import gc 
import sys
import subprocess 
import time
import tkinter as tk
from tkinter import filedialog

# --- PyInstaller パス解決ロジック ---
if getattr(sys, 'frozen', False):
    bundle_dir = sys._MEIPASS
    current_dir = os.path.dirname(sys.executable)
else:
    bundle_dir = os.path.dirname(os.path.abspath(__file__))
    current_dir = bundle_dir

# --- 定数設定 ---
BASE_DB_DIR = os.path.join(current_dir, "chroma_db")
MODELS_DIR = os.path.join(current_dir, "models")
VALID_EXTS = ('.pdf', '.pptx', '.docx', '.doc', '.xlsx', '.xls', '.txt')

# Embeddingモデル設定 (PLamo/E5)
EMBED_MODELS = {
    "Multilingual-E5-Small (軽量・高速)": {
        "id": "intfloat/multilingual-e5-small",
        "dir": "e5_small"
    },
    "PLamo Embedding 1B (高精度・国産)": {
        "id": "pfnet/plamo-embedding-1b",
        "dir": "plamo_1b"
    }
}

DEFAULT_LLM = "gemma3:4b"
DEFAULT_EMBED_LABEL = "Multilingual-E5-Small (軽量・高速)"

# フォルダ生成
for d in [BASE_DB_DIR, MODELS_DIR]:
    if not os.path.exists(d):
        os.makedirs(d)

# 環境変数を固定
os.environ["HUGGINGFACE_HUB_CACHE"] = MODELS_DIR
os.environ["SENTENCE_TRANSFORMERS_HOME"] = MODELS_DIR

# --- インポート (遅延ロード対応) ---
from langchain_community.document_loaders import (
    PyMuPDFLoader, TextLoader, UnstructuredExcelLoader, Docx2txtLoader
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import Ollama
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

# リランカー機能は削除（Ollama専用版）
HAS_RERANKER_LIB = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

st.set_page_config(page_title="Local RAG Pro (PLamo & Rerank)", layout="wide")

# --- 便利関数 ---
def open_in_explorer(file_path):
    if os.path.exists(file_path):
        subprocess.run(['explorer', '/select,', os.path.normpath(file_path)])

def safe_rmtree(path):
    if not os.path.exists(path):
        return True
    gc.collect() 
    time.sleep(0.5)
    for i in range(5): 
        try:
            if os.path.exists(path):
                shutil.rmtree(path)
            return True
        except:
            time.sleep(1)
    return False

# --- セッション状態の初期化 ---
if 'scanning' not in st.session_state: st.session_state.scanning = False
if 'stop_requested' not in st.session_state: st.session_state.stop_requested = False
if 'folder_path' not in st.session_state: st.session_state.folder_path = ""
if 'messages' not in st.session_state: st.session_state.messages = []
if 'last_embed_label' not in st.session_state: st.session_state.last_embed_label = DEFAULT_EMBED_LABEL
if 'scan_percent' not in st.session_state: st.session_state.scan_percent = 0
if 'skipped_files' not in st.session_state: st.session_state.skipped_files = [] 

def select_folder_dialog():
    root = tk.Tk()
    root.withdraw()
    root.attributes('-topmost', True)
    folder_path = filedialog.askdirectory(master=root)
    root.destroy()
    return folder_path

@st.cache_resource
def get_models(llm_model_name, embed_model_id):
    try:
        llm = Ollama(model=llm_model_name)
    except:
        llm = None
    try:
        model_kwargs = {'device': 'cpu', 'trust_remote_code': True}
        embeddings = HuggingFaceEmbeddings(
            model_name=embed_model_id,
            cache_folder=MODELS_DIR,
            model_kwargs=model_kwargs,
            encode_kwargs={'normalize_embeddings': True}
        )
        return llm, embeddings
    except Exception as e:
        st.error(f"モデルロードエラー: {e}")
        return llm, None

# リランカー関数は削除（Ollama専用版）

def get_db_stats(target_path, db_dir):
    if not target_path or not os.path.exists(target_path):
        return 0, 0, 0
    all_files = [os.path.normpath(os.path.join(r, f)) for r, _, fs in os.walk(target_path) for f in fs if f.lower().endswith(VALID_EXTS)]
    total_files = len(all_files)
    if total_files == 0: return 0, 0, 0
    indexed_count = 0
    if os.path.exists(db_dir):
        try:
            db = Chroma(persist_directory=db_dir)
            data = db.get()
            if data and 'metadatas' in data:
                indexed_files_in_db = {os.path.normpath(str(m['source'])).lower() for m in data['metadatas'] if m and 'source' in m}
                indexed_count = sum(1 for f in all_files if f.lower() in indexed_files_in_db)
            db = None
            gc.collect()
        except: pass
    percent = int((indexed_count / total_files) * 100) if total_files > 0 else 0
    return total_files, indexed_count, percent

def load_pptx_lightweight(file_path):
    if not HAS_PPTX:
        raise ImportError("python-pptx がインストールされていません。")
    prs = pptx.Presentation(file_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return [Document(page_content="\n".join(text_content), metadata={"source": file_path})]

def run_scan_and_indexing(target_path, embed_model_id, current_db_dir):
    _, embeddings = get_models(DEFAULT_LLM, embed_model_id)
    if embeddings is None:
        st.session_state.scanning = False
        return False

    db = Chroma(persist_directory=current_db_dir, embedding_function=embeddings)
    indexed_files = set()
    if os.path.exists(current_db_dir):
        try:
            data = db.get()
            if data and 'metadatas' in data:
                indexed_files = {os.path.normpath(str(m['source'])).lower() for m in data['metadatas'] if m and 'source' in m}
        except: pass

    all_files = [os.path.normpath(os.path.join(r, f)) for r, _, fs in os.walk(target_path) for f in fs if f.lower().endswith(VALID_EXTS)]
    new_files = [f for f in all_files if f.lower() not in indexed_files]
    
    st.session_state.skipped_files = [] 
    progress_bar = st.progress(0)
    status_text = st.empty()

    if not new_files:
        st.session_state.scan_percent = 100
        progress_bar.progress(1.0)
        status_text.success("✅ すべて最新の状態です。")
        st.session_state.scanning = False
        return True

    st.session_state.scanning = True
    st.session_state.stop_requested = False
    batch_size = 5 
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    for i in range(0, len(new_files), batch_size):
        if st.session_state.stop_requested:
            st.session_state.scanning = False
            db = None
            gc.collect()
            return False

        current_batch = new_files[i : i + batch_size]
        docs_to_add = []
        
        for idx, f in enumerate(current_batch):
            current_total_idx = i + idx + 1
            st.session_state.scan_percent = int((current_total_idx / len(new_files)) * 100)
            status_text.markdown(f"⏳ **[{current_total_idx}/{len(new_files)}] スキャン中...**\nファイル: `{os.path.basename(f)}`")
            progress_bar.progress(current_total_idx / len(new_files))
            
            try:
                ext = os.path.splitext(f)[1].lower()
                if ext == ".pdf": loader = PyMuPDFLoader(f)
                elif ext == ".pptx": raw_docs = load_pptx_lightweight(f)
                elif ext in [".docx", ".doc"]:
                    loader = Docx2txtLoader(f)
                    raw_docs = loader.load()
                elif ext == ".txt":
                    loader = TextLoader(f, encoding="utf-8")
                    raw_docs = loader.load()
                elif ext in [".xlsx", ".xls"]:
                    loader = UnstructuredExcelLoader(f)
                    raw_docs = loader.load()
                else: continue
                
                splits = text_splitter.split_documents(raw_docs)
                for d in splits:
                    if not d.page_content or not d.page_content.strip(): continue
                    cleaned_meta = {}
                    for k, v in d.metadata.items():
                        if v is not None:
                            cleaned_meta[k] = v if isinstance(v, (str, int, float, bool)) else str(v)
                    d.metadata = cleaned_meta
                    docs_to_add.append(d)
            except Exception as e:
                st.session_state.skipped_files.append(f"{os.path.basename(f)} (Error: {str(e)})")
                continue

        if docs_to_add:
            try:
                db.add_documents(docs_to_add)
            except Exception as e:
                st.session_state.skipped_files.append(f"DB保存エラー: {str(e)}")
        gc.collect()

    status_text.success(f"✅ スキャン完了 ({len(new_files)} 件)")
    st.session_state.scanning = False
    st.session_state.stop_requested = False
    db = None
    gc.collect()
    return True

# --- Ollama関連のヘルパー関数 ---
def check_ollama_running():
    """Ollamaが実行中かチェック"""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        return response.status_code == 200
    except:
        return False

def get_ollama_models():
    """利用可能なOllamaモデルのリストを取得"""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        if response.status_code == 200:
            data = response.json()
            models = [model['name'] for model in data.get('models', [])]
            return models if models else []
        return []
    except:
        return []

# --- UIレイアウト ---
st.sidebar.header("⚙️ システム設定")

# Ollamaステータスチェック
if check_ollama_running():
    st.sidebar.success("✓ Ollama 実行中")
    ollama_models = get_ollama_models()
    if ollama_models:
        selected_ollama_model = st.sidebar.selectbox("LLMモデル:", options=ollama_models, index=0)
        DEFAULT_LLM = selected_ollama_model
    else:
        st.sidebar.warning("モデルがインストールされていません")
        DEFAULT_LLM = "gemma2:2b"
else:
    st.sidebar.error("⚠ Ollama 未起動")
    with st.sidebar.expander("📖 Ollamaのインストール方法", expanded=True):
        st.markdown("""
        1. [Ollama公式サイト](https://ollama.com/download)からダウンロード
        2. インストール後、コマンドプロンプトで:
        ```
        ollama run gemma2:2b
        ```
        3. Ollamaを起動してから再度アクセスしてください
        """)
    DEFAULT_LLM = "gemma2:2b"

st.sidebar.divider()

current_label = st.sidebar.selectbox("Embeddingモデル:", options=list(EMBED_MODELS.keys()), index=0)

if current_label != st.session_state.last_embed_label:
    st.session_state.last_embed_label = current_label
    st.session_state.messages = []
    st.session_state.scanning = False
    st.session_state.stop_requested = False
    st.session_state.scan_percent = 0
    st.rerun()

selected_embed_id = EMBED_MODELS[current_label]["id"]
current_db_dir = os.path.normpath(os.path.join(BASE_DB_DIR, EMBED_MODELS[current_label]["dir"]))

st.sidebar.divider()
st.sidebar.header("📁 データ管理")
if st.sidebar.button("📁 フォルダを選択"):
    path = select_folder_dialog()
    if path: st.session_state.folder_path = os.path.normpath(path)

total_f, indexed_f, db_percent = get_db_stats(st.session_state.folder_path, current_db_dir)

if st.session_state.folder_path:
    st.sidebar.write(f"対象フォルダ: `{st.session_state.folder_path}`")
    st.sidebar.info(f"📊 登録状況: {indexed_f} / {total_f} ファイル ({db_percent}%)")
    
    if not st.session_state.scanning:
        if st.sidebar.button("⚡ スキャン開始/再開"):
            run_scan_and_indexing(st.session_state.folder_path, selected_embed_id, current_db_dir)
            st.rerun()
    else:
        st.sidebar.warning(f"⏳ スキャン実行中... ({st.session_state.scan_percent}%)")
        if st.sidebar.button("🛑 スキャンを中断"):
            st.session_state.stop_requested = True
            st.rerun()

if st.session_state.skipped_files:
    with st.sidebar.expander("⚠️ スキップされたファイル", expanded=False):
        for skipped in st.session_state.skipped_files:
            st.caption(skipped)

if st.sidebar.button("🗑️ このDBをクリア"):
    if safe_rmtree(current_db_dir):
        st.sidebar.success("DBをクリアしました。")
        st.session_state.messages = []
        st.rerun()

st.title("📂 Local RAG System Pro")

col1, col2, col3 = st.columns(3)
with col1: st.metric("🧠 LLM", DEFAULT_LLM)
with col2: st.metric("🧬 Embedding", current_label.split(" ")[0])
with col3:
    status_label = "構築済み" if db_percent == 100 else "構築中" if db_percent > 0 else "未構築"
    st.metric("💾 DB状態", f"{status_label} ({db_percent}%)")

st.divider()

for i, m in enumerate(st.session_state.messages):
    with st.chat_message(m["role"]):
        st.markdown(m["content"])
        if m["role"] == "assistant" and m.get("sources"):
            st.caption("📂 **参照元のファイル:**")
            source_cols = st.columns(min(len(m["sources"]), 3))
            for idx, src in enumerate(m["sources"]):
                with source_cols[idx % 3]:
                    if st.button(f"📄 {os.path.basename(src)}", key=f"hist_{i}_{idx}"):
                        open_in_explorer(src)

if prompt_input := st.chat_input("質問を入力してください..."):
    st.session_state.messages.append({"role": "user", "content": prompt_input})
    st.rerun()

if st.session_state.messages and st.session_state.messages[-1]["role"] == "user":
    last_prompt = st.session_state.messages[-1]["content"]
    with st.chat_message("assistant"):
        llm, embeddings = get_models(DEFAULT_LLM, selected_embed_id)
        if llm and embeddings:
            db = Chroma(persist_directory=current_db_dir, embedding_function=embeddings)
            
            with st.spinner("情報を検索・分析中..."):
                # 通常検索（リランカーは削除）
                retriever = db.as_retriever(search_kwargs={"k": 3})
                docs = retriever.invoke(last_prompt)

                context_text = "\n\n".join([f"【出典: {os.path.basename(d.metadata.get('source', ''))}】\n{d.page_content}" for d in docs])
                
                source_list = []
                for d in docs:
                    s = d.metadata.get("source")
                    if s and s not in source_list: source_list.append(s)
                
                prompt_tpl = ChatPromptTemplate.from_template(
                    "以下の資料を参考に、質問に対して誠実に日本語で回答してください。資料にない情報は無理に答えないでください。\n\n資料:\n{context}\n\n質問: {question}"
                )
                chain = prompt_tpl | llm | StrOutputParser()
                
                full_response = st.write_stream(chain.stream({"context": context_text, "question": last_prompt}))
                
                if source_list:
                    st.caption("📂 **参照元のファイル:**")
                    cols = st.columns(min(len(source_list), 3))
                    for idx, src in enumerate(source_list):
                        with cols[idx % 3]:
                            if st.button(f"📄 {os.path.basename(src)}", key=f"cur_{idx}"):
                                open_in_explorer(src)
                
                st.session_state.messages.append({"role": "assistant", "content": full_response, "sources": source_list})
        else:
            st.error("モデルのロードに失敗しました。Ollamaが起動しているか確認してください。")

if __name__ == "__main__":
    pass